"""批量转换脚本 - 多线程 + 优先级调度 + 保留目录树

架构:
  - local_pool (独享线程): 纯本地解析文档 (docx/xlsx/pptx/pdf 字多图少)
  - api_pool   (N线程):   需要调 Gemini/OpenAI 的文档 (图片 / 图多字少)

优先级:
  0 = 纯本地 (docx/xlsx/pptx/pdf/doc/ppt/xls 等)
  1 = 图片文件 (jpg/png/... 一定调 API)

用法:
    python -m omniparser.batch_convert \\
        --input  /path/to/source/ \\
        --output /path/to/target/ \\
        [--format both] [--verbose] [--workers 4]
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from .config import Config
from .pipeline import Pipeline
from .utils import setup_logging, ensure_dir, collect_files

# 图片扩展名 — 一定需要 API
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}

# 可能含内嵌图片的文档格式
DOC_WITH_IMAGES = {".pdf", ".docx", ".pptx"}

# 直接复制的纯文本格式（不需要解析）
DIRECT_COPY_EXTENSIONS = {".txt"}

# 跳过的归档/多媒体格式
SKIP_EXTENSIONS = {".zip", ".rar", ".7z", ".tar", ".gz", ".mp4", ".avi", ".mov", ".eml"}

_prescan_logger = logging.getLogger("omniparser.prescan")


def _is_image_file(file_path: Path) -> bool:
    return file_path.suffix.lower() in IMAGE_EXTENSIONS


def _prescan_needs_api(file_path: Path) -> bool:
    """快速预扫描文档，判断是否图多字少（需要 API）

    只对 PDF / DOCX / PPTX 做轻量扫描，其他格式直接返回 False
    """
    ext = file_path.suffix.lower()
    if ext not in DOC_WITH_IMAGES:
        return False

    from .parsers.image_describer import TEXT_PER_IMAGE_THRESHOLD

    try:
        if ext == ".pdf":
            return _prescan_pdf(file_path, TEXT_PER_IMAGE_THRESHOLD)
        elif ext == ".docx":
            return _prescan_docx(file_path, TEXT_PER_IMAGE_THRESHOLD)
        elif ext == ".pptx":
            return _prescan_pptx(file_path, TEXT_PER_IMAGE_THRESHOLD)
    except Exception as e:
        _prescan_logger.debug("预扫描失败 %s: %s，归入本地队列", file_path.name, e)
    return False


def _prescan_pdf(file_path: Path, threshold: int) -> bool:
    import fitz

    doc = fitz.open(str(file_path))
    text_chars = 0
    image_count = 0
    for page in doc:
        text_chars += len(page.get_text("text").strip())
        image_count += len(page.get_images(full=False))
    doc.close()
    if image_count == 0:
        return False
    return (text_chars / image_count) < threshold


def _prescan_docx(file_path: Path, threshold: int) -> bool:
    from docx import Document as DocxDocument

    docx = DocxDocument(str(file_path))
    text_chars = 0
    image_count = 0
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for element in docx.element.body:
        tag = element.tag.split("}")[-1]
        if tag == "p":
            from docx.text.paragraph import Paragraph

            text_chars += len(Paragraph(element, docx).text.strip())
            image_count += len(element.findall(".//a:blip", ns))
        elif tag == "tbl":
            from docx.table import Table

            for row in Table(element, docx).rows:
                for cell in row.cells:
                    text_chars += len(cell.text.strip())
    if image_count == 0:
        return False
    return (text_chars / image_count) < threshold


def _prescan_pptx(file_path: Path, threshold: int) -> bool:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(file_path))
    text_chars = 0
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_chars += len(para.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
    if image_count == 0:
        return False
    return (text_chars / image_count) < threshold


def _is_skipped(out_subdir: Path, stem: str, output_format: str) -> bool:
    """检查输出文件是否已存在"""
    md_exists = (out_subdir / f"{stem}.md").exists()
    json_exists = (out_subdir / f"{stem}.json").exists()
    if output_format == "both":
        return md_exists and json_exists
    elif output_format == "markdown":
        return md_exists
    elif output_format == "json":
        return json_exists
    return False


def _copy_txt(
    file_path: Path, rel_path: Path, out_subdir: Path, stem: str, output_format: str
) -> dict:
    """直接复制 .txt 文件内容到输出（不需要解析）"""
    logger = logging.getLogger("omniparser.batch")
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace").strip()
        if not text:
            text = f"> (空文件: {file_path.name})"

        if output_format in ("markdown", "both"):
            md_path = out_subdir / f"{stem}.md"
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(f"<!-- source: {rel_path} -->\n\n")
                f.write(text)
                f.write("\n")

        if output_format in ("json", "both"):
            json_path = out_subdir / f"{stem}.json"
            data = {
                "source": str(rel_path),
                "documents": [
                    {
                        "content": text,
                        "content_type": "paragraph",
                        "metadata": {"direct_copy": True},
                    }
                ],
                "chunks": [{"text": text}],
                "error": None,
            }
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)

        logger.info("COPY: %s (txt直接复制)", rel_path)
        return {"status": "success", "file": str(rel_path)}
    except Exception as e:
        logger.error("COPY FAILED: %s - %s", rel_path, e)
        return {"status": "failed", "file": str(rel_path), "error": str(e)}


def _process_one(
    file_path: Path,
    input_dir: Path,
    output_dir: Path,
    pipeline: Pipeline,
    output_format: str,
) -> dict:
    """处理单个文件，返回结果 dict"""
    rel_path = file_path.relative_to(input_dir)
    out_subdir = output_dir / rel_path.parent
    ensure_dir(out_subdir)
    stem = rel_path.stem

    logger = logging.getLogger("omniparser.batch")

    # --- .txt 直接复制 ---
    if file_path.suffix.lower() in DIRECT_COPY_EXTENSIONS:
        return _copy_txt(file_path, rel_path, out_subdir, stem, output_format)

    try:
        result = pipeline.parse_file(file_path)
    except Exception as e:
        logger.error("EXCEPTION: %s - %s", rel_path, e)
        return {"status": "failed", "file": str(rel_path), "error": str(e)}

    if not result.success:
        logger.warning("FAILED: %s - %s", rel_path, result.error)
        return {"status": "failed", "file": str(rel_path), "error": result.error}

    # 写入输出
    try:
        if output_format in ("json", "both"):
            json_path = out_subdir / f"{stem}.json"
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(result.to_dict(), f, ensure_ascii=False, indent=2)

        if output_format in ("markdown", "both"):
            md_path = out_subdir / f"{stem}.md"
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(f"<!-- source: {rel_path} -->\n\n")
                for doc in result.documents:
                    f.write(doc.content)
                    f.write("\n\n")

        docs_count = len(result.documents)
        chunks_count = len(result.chunks)
        logger.info("OK: %s (%d docs, %d chunks)", rel_path, docs_count, chunks_count)
        return {"status": "success", "file": str(rel_path)}

    except Exception as e:
        logger.error("Write error: %s - %s", rel_path, e)
        return {"status": "failed", "file": str(rel_path), "error": f"Write error: {e}"}


def batch_convert(
    input_dir: Path,
    output_dir: Path,
    config: Config,
    output_format: str = "both",
    api_workers: int = 8,
) -> dict:
    """多线程批量转换

    Args:
        input_dir: 源目录
        output_dir: 目标目录
        config: 配置
        output_format: 输出格式
        api_workers: API 线程池大小

    Returns:
        统计信息 dict
    """
    logger = logging.getLogger("omniparser.batch")

    input_dir = Path(input_dir).resolve()
    output_dir = Path(output_dir).resolve()

    pipeline = Pipeline(config)

    # 收集所有文件：pipeline 支持的 + .txt 直接复制
    all_extensions = pipeline.supported_extensions | DIRECT_COPY_EXTENSIONS
    all_files = collect_files(input_dir, recursive=True, extensions=all_extensions)
    total = len(all_files)

    # --- 分类 & 跳过 ---
    local_queue: list[Path] = []  # 纯本地解析（字多图少 / 无图 / txt直接复制）
    api_queue: list[Path] = []  # 需要 API（图片文件 + 图多字少文档）
    skipped = 0

    logger.info("预扫描文件，分类中...")
    for fp in all_files:
        rel_path = fp.relative_to(input_dir)
        out_subdir = output_dir / rel_path.parent
        stem = rel_path.stem
        if _is_skipped(out_subdir, stem, output_format):
            skipped += 1
            continue
        if fp.suffix.lower() in DIRECT_COPY_EXTENSIONS:
            local_queue.append(fp)  # .txt 走本地队列直接复制
        elif _is_image_file(fp):
            api_queue.append(fp)
        elif _prescan_needs_api(fp):
            _prescan_logger.info("图多字少 → API队列: %s", rel_path.name)
            api_queue.append(fp)
        else:
            local_queue.append(fp)

    need_process = len(local_queue) + len(api_queue)
    logger.info(
        "共 %d 文件: %d 本地解析, %d 需API, %d 已跳过",
        total,
        len(local_queue),
        len(api_queue),
        skipped,
    )

    stats = {
        "total": total,
        "success": 0,
        "failed": 0,
        "skipped": skipped,
        "errors": [],
    }
    stats_lock = threading.Lock()

    # --- 进度条 ---
    try:
        from tqdm import tqdm

        pbar = tqdm(
            total=need_process,
            desc="转换进度",
            unit="file",
            bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}, {rate_fmt}]",
        )
    except ImportError:
        pbar = None

    def _update_stats(result: dict):
        with stats_lock:
            if result["status"] == "success":
                stats["success"] += 1
            else:
                stats["failed"] += 1
                stats["errors"].append(result)
        if pbar:
            pbar.update(1)
            short = Path(result["file"]).name
            tag = "✓" if result["status"] == "success" else "✗"
            pbar.set_postfix_str(f"{tag} {short}")

    start_time = time.time()

    # --- 并行执行 ---
    # 1. 本地线程池 (1 线程独享，不抢 API 资源)
    # 2. API 线程池 (N 线程并发调 Gemini)
    local_futures = []
    api_futures = []

    local_pool = ThreadPoolExecutor(max_workers=1, thread_name_prefix="local")
    api_pool = ThreadPoolExecutor(max_workers=api_workers, thread_name_prefix="api")

    # 优先提交本地任务（它们不等 API，跑得快）
    for fp in local_queue:
        fut = local_pool.submit(
            _process_one, fp, input_dir, output_dir, pipeline, output_format
        )
        local_futures.append(fut)

    # 提交 API 任务
    for fp in api_queue:
        fut = api_pool.submit(
            _process_one, fp, input_dir, output_dir, pipeline, output_format
        )
        api_futures.append(fut)

    # 等待所有完成
    all_futures = local_futures + api_futures
    for fut in as_completed(all_futures):
        try:
            result = fut.result()
            _update_stats(result)
        except Exception as e:
            logger.error("Unexpected thread error: %s", e)
            with stats_lock:
                stats["failed"] += 1
                stats["errors"].append({"file": "unknown", "error": str(e)})
            if pbar:
                pbar.update(1)

    local_pool.shutdown(wait=True)
    api_pool.shutdown(wait=True)

    if pbar:
        pbar.close()

    elapsed = time.time() - start_time
    stats["elapsed_seconds"] = round(elapsed, 1)

    return stats


def main():
    parser = argparse.ArgumentParser(
        description="OmniParser 批量转换 - 多线程 + 优先级调度",
    )
    parser.add_argument("-i", "--input", required=True, help="源文件目录")
    parser.add_argument("-o", "--output", required=True, help="输出目录")
    parser.add_argument(
        "-f",
        "--format",
        choices=["json", "markdown", "both"],
        default="both",
        help="输出格式 (默认: both)",
    )
    parser.add_argument("-c", "--config", default="config.yaml", help="配置文件路径")
    parser.add_argument("-v", "--verbose", action="store_true", help="详细日志")
    parser.add_argument(
        "-w",
        "--workers",
        type=int,
        default=8,
        help="API 并发线程数 (默认: 8)",
    )

    args = parser.parse_args()

    # 自动加载 .env 文件
    _load_dotenv()

    level = "DEBUG" if args.verbose else "INFO"
    setup_logging(level)
    for noisy in ("httpcore", "httpx", "urllib3", "google_genai", "google.auth"):
        logging.getLogger(noisy).setLevel(logging.WARNING)
    logger = logging.getLogger("omniparser.batch")

    input_dir = Path(args.input)
    output_dir = Path(args.output)

    if not input_dir.exists():
        logger.error("Input directory does not exist: %s", input_dir)
        sys.exit(1)

    config = Config.from_yaml(args.config)

    print(f"\n{'='*60}")
    print(f"OmniParser 批量转换 (多线程)")
    print(f"  源目录:     {input_dir}")
    print(f"  输出:       {output_dir}")
    print(f"  格式:       {args.format}")
    print(f"  本地线程:   1 (独享)")
    print(f"  API 线程:   {args.workers}")
    print(f"{'='*60}\n")

    stats = batch_convert(
        input_dir,
        output_dir,
        config,
        args.format,
        api_workers=args.workers,
    )

    print(f"\n{'='*60}")
    print(f"转换完成!")
    print(f"  总文件: {stats['total']}")
    print(f"  成功:   {stats['success']}")
    print(f"  失败:   {stats['failed']}")
    print(f"  跳过:   {stats['skipped']}")
    print(f"  耗时:   {stats['elapsed_seconds']}s")
    print(f"{'='*60}")

    if stats["errors"]:
        error_log = output_dir / "_errors.json"
        with open(error_log, "w", encoding="utf-8") as f:
            json.dump(stats["errors"], f, ensure_ascii=False, indent=2)
        print(f"\n失败详情已写入: {error_log}")

    stats_path = output_dir / "_stats.json"
    with open(stats_path, "w", encoding="utf-8") as f:
        json.dump(stats, f, ensure_ascii=False, indent=2)

    if stats["failed"] > 0:
        sys.exit(1)


def _load_dotenv():
    """从项目根目录 .env 文件加载环境变量（不覆盖已有的）"""
    env_file = Path(__file__).resolve().parent.parent / ".env"
    if not env_file.exists():
        return
    import os

    for line in env_file.read_text().splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, _, value = line.partition("=")
        key, value = key.strip(), value.strip()
        if key and value and key not in os.environ:
            os.environ[key] = value


if __name__ == "__main__":
    main()
