"""PPTX / PPT 文件解析器"""

from __future__ import annotations

import logging
import re
import struct
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ..models import Document, ContentType
from ..config import Config
from .base import BaseParser

logger = logging.getLogger("omniparser.parsers.pptx")


class PptxParser(BaseParser):
    """解析 .pptx 和 .ppt 文件

    - .pptx 直接用 python-pptx 解析
    - .ppt 降级为 olefile 纯 Python 提取文本
    """

    supported_extensions = {".pptx", ".ppt"}

    def parse(self, file_path: Path) -> list[Document]:
        if file_path.suffix.lower() == ".ppt":
            return self._parse_ppt(file_path)
        return self._parse_pptx(file_path)

    def _parse_ppt(self, file_path: Path) -> list[Document]:
        """用 olefile 从 .ppt OLE 容器中提取文本"""
        try:
            import olefile
        except ImportError:
            raise ImportError(
                "olefile not installed. Install with: pip install olefile"
            )

        source = str(file_path)

        if not olefile.isOleFile(str(file_path)):
            raise ValueError(f"Not a valid OLE file: {file_path}")

        ole = olefile.OleFileIO(str(file_path))
        documents: list[Document] = []

        try:
            # PPT 文本存储在 "PowerPoint Document" stream 中
            if ole.exists("PowerPoint Document"):
                ppt_stream = ole.openstream("PowerPoint Document").read()
                text = self._extract_ppt_text(ppt_stream)
            else:
                # 尝试从所有 stream 中提取
                text = self._extract_text_from_ole_streams(ole)

            if text.strip():
                # 尝试按幻灯片分隔（简化：按连续空行分段）
                sections = re.split(r"\n{3,}", text.strip())
                for idx, section in enumerate(sections, start=1):
                    section = section.strip()
                    if not section:
                        continue
                    lines = section.split("\n")
                    # 第一行作为标题
                    title = lines[0].strip() if lines else ""
                    if title:
                        md = f"## Slide {idx}: {title}"
                        if len(lines) > 1:
                            body = "\n".join(
                                f"- {l.strip()}" for l in lines[1:] if l.strip()
                            )
                            md += f"\n\n{body}"
                    else:
                        md = f"## Slide {idx}\n\n{section}"

                    documents.append(
                        Document(
                            source=source,
                            content=md,
                            content_type=ContentType.PARAGRAPH,
                            page=idx,
                            metadata={"parse_method": "olefile", "slide_number": idx},
                        )
                    )
            else:
                documents.append(
                    Document(
                        source=source,
                        content=f"> (无法从 .ppt 文件提取文本: {file_path.name})",
                        content_type=ContentType.PARAGRAPH,
                        metadata={"parse_method": "olefile", "partial": True},
                    )
                )
        finally:
            ole.close()

        return documents

    def _extract_ppt_text(self, stream_data: bytes) -> str:
        """从 PowerPoint Document stream 提取文本

        PPT 二进制格式中的文本记录类型:
        - TextCharsAtom (0x0FA0): UTF-16LE 文本
        - TextBytesAtom (0x0FA8): ASCII/CP1252 文本
        """
        texts = []
        offset = 0
        data_len = len(stream_data)

        while offset < data_len - 8:
            try:
                # 读取 record header: recVer(4bit) + recInstance(12bit) + recType(16bit) + recLen(32bit)
                rec_ver_inst = struct.unpack_from("<H", stream_data, offset)[0]
                rec_type = struct.unpack_from("<H", stream_data, offset + 2)[0]
                rec_len = struct.unpack_from("<I", stream_data, offset + 4)[0]

                if rec_type == 0x0FA0:  # TextCharsAtom - UTF-16LE
                    text_data = stream_data[offset + 8 : offset + 8 + rec_len]
                    try:
                        text = text_data.decode("utf-16-le", errors="ignore")
                        text = text.strip()
                        if text and len(text) > 1:
                            texts.append(text)
                    except Exception:
                        pass
                    offset += 8 + rec_len

                elif rec_type == 0x0FA8:  # TextBytesAtom - ASCII/CP1252
                    text_data = stream_data[offset + 8 : offset + 8 + rec_len]
                    try:
                        text = text_data.decode("cp1252", errors="ignore")
                        text = text.strip()
                        if text and len(text) > 1:
                            texts.append(text)
                    except Exception:
                        pass
                    offset += 8 + rec_len

                else:
                    # 非文本记录，跳过 header
                    # 容器记录 (recVer == 0xF) 的 data 包含子记录，只跳过 header
                    rec_ver = rec_ver_inst & 0x0F
                    if rec_ver == 0x0F:
                        offset += 8  # 进入容器
                    else:
                        offset += 8 + rec_len  # 跳过整个记录

            except struct.error:
                break

        return "\n\n".join(texts)

    def _extract_text_from_ole_streams(self, ole) -> str:
        """从 OLE 的所有 stream 中尝试提取文本"""
        best_text = ""
        for stream_path in ole.listdir():
            try:
                data = ole.openstream(stream_path).read()
                text = data.decode("utf-16-le", errors="ignore")
                text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
                readable = sum(1 for c in text if c.isalnum() or c in "，。！？")
                if readable > len(best_text):
                    best_text = text.strip()
            except Exception:
                continue
        return best_text

    def _parse_pptx(self, file_path: Path) -> list[Document]:
        """解析 .pptx，内嵌图片按字图比例决定是否解析"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from .image_describer import (
            should_parse_images,
            describe_images_batch,
            ImageItem,
        )

        source = str(file_path)
        documents: list[Document] = []

        prs = Presentation(str(file_path))

        # --- 第一遍：扫描统计文本量和图片数 ---
        total_text_chars = 0
        total_image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        total_text_chars += len(para.text.strip())
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total_image_count += 1

        parse_images = should_parse_images(total_text_chars, total_image_count)

        # --- 收集所有图片（先收集后合批） ---
        image_items: list[ImageItem] = []
        # 记录每张图片应该插入到哪个 slide 的 body_texts 中
        # key: image 在 image_items 中的索引
        # value: (slide_idx_0based, placeholder_string)
        image_placeholders: dict[int, tuple[int, str]] = {}

        # --- 第二遍：提取文本 + 收集图片占位 ---
        slide_data: list[tuple[str, list[str]]] = []  # [(title_md, body_texts), ...]

        for slide_idx, slide in enumerate(prs.slides, start=1):
            # 提取标题
            title = self._extract_title(slide)
            if title:
                title_md = f"## Slide {slide_idx}: {title}"
            else:
                title_md = f"## Slide {slide_idx}"

            body_texts: list[str] = [""]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            level = paragraph.level or 0
                            if level > 0:
                                indent = "  " * level
                                body_texts.append(f"{indent}- {text}")
                            else:
                                body_texts.append(f"- {text}")

                if shape.has_table:
                    table_md = self._extract_table(shape.table)
                    if table_md:
                        body_texts.append("")
                        body_texts.append(table_md)

                # 收集图片
                if parse_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        img_ct = shape.image.content_type or "image/png"
                        img_idx = len(image_items)
                        image_items.append(
                            ImageItem(
                                image_bytes=img_blob,
                                mime_type=img_ct,
                                context=f"PPT第{slide_idx}页",
                            )
                        )
                        placeholder = f"__IMG_PLACEHOLDER_{img_idx}__"
                        body_texts.append("")
                        body_texts.append(placeholder)
                        image_placeholders[img_idx] = (slide_idx - 1, placeholder)
                    except Exception as e:
                        logger.warning("Failed to extract PPTX image: %s", e)

            slide_data.append((title_md, body_texts))

        # --- 合批描述图片 ---
        if image_items:
            describe_images_batch(image_items, self.config)

        # --- 组装最终文档 ---
        for slide_idx_0, (title_md, body_texts) in enumerate(slide_data):
            md_parts = [title_md]

            for line in body_texts:
                # 替换图片占位符
                replaced = False
                for img_idx, (s_idx, placeholder) in image_placeholders.items():
                    if line == placeholder:
                        desc = image_items[img_idx].description
                        if desc:
                            md_parts.append("")
                            md_parts.append(desc)
                        replaced = True
                        break
                if not replaced:
                    md_parts.append(line)

            content = "\n".join(md_parts)
            if content.strip():
                documents.append(
                    Document(
                        source=source,
                        content=content,
                        content_type=ContentType.PARAGRAPH,
                        page=slide_idx_0 + 1,
                        metadata={"slide_number": slide_idx_0 + 1},
                    )
                )

        return documents

    def _extract_title(self, slide) -> str | None:
        """提取幻灯片标题"""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text_frame.text.strip() or None

        # 备用：查找第一个占位符类型为标题的
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # 标题占位符
                if shape.has_text_frame:
                    return shape.text_frame.text.strip() or None
        return None

    def _extract_table(self, table) -> str:
        """将 PPT 表格转为 Markdown"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        lines = []
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        for row in rows[1:]:
            while len(row) < len(rows[0]):
                row.append("")
            lines.append("| " + " | ".join(row[: len(rows[0])]) + " |")

        return "\n".join(lines)
