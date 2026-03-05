"""OmniParser 核心模块测试"""

import json
import tempfile
from pathlib import Path

import pytest

from omniparser.config import Config
from omniparser.models import Document, Chunk, ParseResult, ContentType
from omniparser.pipeline import Pipeline
from omniparser.cache import FileCache
from omniparser.utils import compute_file_hash, collect_files
from omniparser.postprocessors.chunker import Chunker
from omniparser.postprocessors.metadata import MetadataExtractor


# ─── Models ───


class TestDocument:
    def test_to_dict(self):
        doc = Document(
            source="test.pdf",
            content="Hello world",
            content_type=ContentType.PARAGRAPH,
            page=1,
        )
        d = doc.to_dict()
        assert d["source"] == "test.pdf"
        assert d["content_type"] == "paragraph"
        assert d["page"] == 1
        assert "sheet" not in d  # None 值应该被过滤

    def test_to_json(self):
        doc = Document(source="test.pdf", content="你好世界")
        j = doc.to_json()
        data = json.loads(j)
        assert data["content"] == "你好世界"


class TestParseResult:
    def test_success(self):
        result = ParseResult(source="test.pdf")
        assert result.success is True

    def test_failure(self):
        result = ParseResult(source="test.pdf", error="Something went wrong")
        assert result.success is False
        assert result.error == "Something went wrong"


# ─── Utils ───


class TestUtils:
    def test_compute_file_hash(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("hello")
        h = compute_file_hash(f)
        assert h.startswith("sha256:")
        assert len(h) > 10

    def test_collect_files_single_file(self, tmp_path):
        f = tmp_path / "test.pdf"
        f.touch()
        files = collect_files(f)
        assert len(files) == 1

    def test_collect_files_dir(self, tmp_path):
        (tmp_path / "a.pdf").touch()
        (tmp_path / "b.docx").touch()
        (tmp_path / "c.txt").touch()
        files = collect_files(tmp_path, extensions={".pdf", ".docx"})
        assert len(files) == 2

    def test_collect_files_recursive(self, tmp_path):
        sub = tmp_path / "sub"
        sub.mkdir()
        (tmp_path / "a.pdf").touch()
        (sub / "b.pdf").touch()
        files = collect_files(tmp_path, recursive=True, extensions={".pdf"})
        assert len(files) == 2


# ─── Cache ───


class TestCache:
    def test_cache_miss(self, tmp_path):
        from omniparser.config import CacheConfig

        cache = FileCache(CacheConfig(enabled=True, dir=str(tmp_path / "cache")))
        f = tmp_path / "test.txt"
        f.write_text("hello")
        assert cache.get(f) is None

    def test_cache_put_and_get(self, tmp_path):
        from omniparser.config import CacheConfig

        cache = FileCache(CacheConfig(enabled=True, dir=str(tmp_path / "cache")))
        f = tmp_path / "test.txt"
        f.write_text("hello")

        result = ParseResult(
            source=str(f),
            documents=[
                Document(
                    source=str(f), content="hello", content_type=ContentType.PARAGRAPH
                )
            ],
            file_hash=compute_file_hash(f),
        )
        cache.put(result)

        # 应该命中缓存
        cached = cache.get(f)
        assert cached is not None
        assert cached.source == str(f)
        assert len(cached.documents) == 1

    def test_cache_disabled(self, tmp_path):
        from omniparser.config import CacheConfig

        cache = FileCache(CacheConfig(enabled=False, dir=str(tmp_path / "cache")))
        f = tmp_path / "test.txt"
        f.write_text("hello")
        assert cache.get(f) is None

    def test_cache_clear(self, tmp_path):
        from omniparser.config import CacheConfig

        cache_dir = tmp_path / "cache"
        cache = FileCache(CacheConfig(enabled=True, dir=str(cache_dir)))
        # 写入一个假的缓存
        (cache_dir / "test.json").write_text("{}")
        assert cache.clear() == 1


# ─── Chunker ───


class TestChunker:
    def test_chunk_by_heading(self):
        from omniparser.config import ChunkingConfig

        chunker = Chunker(ChunkingConfig(strategy="heading", max_tokens=512))
        docs = [
            Document(
                source="t.md",
                content="# Title\n\nParagraph 1",
                content_type=ContentType.HEADING,
            ),
            Document(
                source="t.md",
                content="## Section\n\nParagraph 2",
                content_type=ContentType.HEADING,
            ),
        ]
        chunks = chunker.chunk_documents(docs)
        assert len(chunks) >= 2

    def test_chunk_by_tokens(self):
        from omniparser.config import ChunkingConfig

        chunker = Chunker(
            ChunkingConfig(strategy="fixed_token", max_tokens=20, overlap=5)
        )
        docs = [Document(source="t.md", content="A" * 200)]
        chunks = chunker.chunk_documents(docs)
        assert len(chunks) > 1


# ─── MetadataExtractor ───


class TestMetadataExtractor:
    def test_enrich(self, tmp_path):
        f = tmp_path / "test.pdf"
        f.write_text("hello")
        doc = Document(source=str(f), content="hello")
        extractor = MetadataExtractor()
        extractor.enrich(doc, f)
        assert doc.metadata["file_name"] == "test.pdf"
        assert doc.metadata["file_ext"] == ".pdf"
        assert doc.metadata["file_size"] == 5


# ─── Pipeline routing ───


class TestPipeline:
    def test_supported_extensions(self):
        config = Config.default()
        pipeline = Pipeline(config)
        exts = pipeline.supported_extensions
        assert ".pdf" in exts
        assert ".docx" in exts
        assert ".xlsx" in exts
        assert ".pptx" in exts
        assert ".png" in exts

    def test_unsupported_format(self, tmp_path):
        f = tmp_path / "test.xyz"
        f.write_text("hello")
        config = Config.default()
        pipeline = Pipeline(config)
        result = pipeline.parse_file(f)
        assert not result.success
        assert "Unsupported" in result.error


# ─── DOCX Parser ───


class TestDocxParser:
    def test_parse_docx(self, tmp_path):
        """创建一个简单的 DOCX 并解析"""
        from docx import Document as DocxDoc

        docx_path = tmp_path / "test.docx"
        doc = DocxDoc()
        doc.add_heading("测试标题", level=1)
        doc.add_paragraph("这是一个测试段落。")
        doc.add_heading("二级标题", level=2)
        doc.add_paragraph("另一个段落。")

        # 添加表格
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "列A"
        table.cell(0, 1).text = "列B"
        table.cell(1, 0).text = "数据1"
        table.cell(1, 1).text = "数据2"

        doc.save(str(docx_path))

        # 解析
        config = Config.default()
        pipeline = Pipeline(config)
        result = pipeline.parse_file(docx_path)

        assert result.success
        assert len(result.documents) > 0

        # 验证标题
        headings = [
            d for d in result.documents if d.content_type == ContentType.HEADING
        ]
        assert any("测试标题" in d.content for d in headings)

        # 验证表格
        tables = [d for d in result.documents if d.content_type == ContentType.TABLE]
        assert len(tables) > 0
        assert "列A" in tables[0].content


# ─── XLSX Parser ───


class TestXlsxParser:
    def test_parse_xlsx(self, tmp_path):
        """创建一个简单的 XLSX 并解析"""
        import pandas as pd

        xlsx_path = tmp_path / "test.xlsx"
        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [30, 25]})
        df.to_excel(str(xlsx_path), index=False, sheet_name="Sheet1")

        config = Config.default()
        pipeline = Pipeline(config)
        result = pipeline.parse_file(xlsx_path)

        assert result.success
        assert len(result.documents) == 1
        assert "Alice" in result.documents[0].content
        assert "Bob" in result.documents[0].content
        assert result.documents[0].sheet == "Sheet1"


# ─── PPTX Parser ───


class TestPptxParser:
    def test_parse_pptx(self, tmp_path):
        """创建一个简单的 PPTX 并解析"""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "测试幻灯片"
        slide.placeholders[1].text = "幻灯片内容"
        prs.save(str(pptx_path))

        config = Config.default()
        pipeline = Pipeline(config)
        result = pipeline.parse_file(pptx_path)

        assert result.success
        assert len(result.documents) > 0
        assert any("测试幻灯片" in d.content for d in result.documents)


# ─── PDF Parser ───


class TestPdfParser:
    def test_parse_text_pdf(self, tmp_path):
        """创建一个包含文本的 PDF 并解析"""
        import fitz

        pdf_path = tmp_path / "test.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello World PDF Test", fontsize=12)
        page.insert_text((72, 100), "This is a test paragraph.", fontsize=10)
        doc.save(str(pdf_path))
        doc.close()

        config = Config.default()
        pipeline = Pipeline(config)
        result = pipeline.parse_file(pdf_path)

        assert result.success
        assert len(result.documents) > 0
        all_content = " ".join(d.content for d in result.documents)
        assert "Hello World" in all_content
